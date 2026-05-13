import asyncio
from datetime import datetime, timezone
from functools import lru_cache
import logging
from pathlib import Path
import re
import tempfile
import time
import uuid

import httpx
from bson import ObjectId
import docx2txt
from langchain_core.documents import Document
from motor.motor_asyncio import AsyncIOMotorClient, AsyncIOMotorDatabase
from pypdf import PdfReader
from pptx import Presentation
from rq import get_current_job
from supabase import create_client

from .config import settings
from .gemini import embed_documents_with_fallback, has_gemini_api_keys

try:
    from llama_cloud_services import LlamaParse

    LLAMAPARSE_AVAILABLE = True
except Exception:
    LlamaParse = None  # type: ignore[assignment]
    LLAMAPARSE_AVAILABLE = False


logger = logging.getLogger(__name__)
LLAMAPARSE_CONTENT_TYPES = {"pdf", "docx", "pptx"}
TOKEN_PATTERN = re.compile(r"\S+")
MARKDOWN_HEADING_PATTERN = re.compile(r"^\s{0,3}#{1,6}\s+.+$")
NUMBERED_HEADING_PATTERN = re.compile(r"^\s*\d+(?:\.\d+){0,3}[\)\.\-:]?\s+\S+")
BULLET_LINE_PATTERN = re.compile(r"^\s*(?:[-*+]|(?:\d+[\.\)]))\s+\S+")


def normalize_chunk_text(text: str) -> str:
    normalized = text.replace("\u00a0", " ")
    normalized = re.sub(r"(\w)-\n(\w)", r"\1\2", normalized)
    normalized = re.sub(r"[ \t]+\n", "\n", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    normalized = re.sub(r"[ \t]{2,}", " ", normalized)
    return normalized.strip()


def tokenize_with_spans(text: str) -> list[tuple[int, int]]:
    return [(match.start(), match.end()) for match in TOKEN_PATTERN.finditer(text)]


def count_tokens(text: str) -> int:
    return len(TOKEN_PATTERN.findall(text))


def split_text_with_token_windows(
    text: str, chunk_tokens: int, overlap_tokens: int
) -> list[str]:
    spans = tokenize_with_spans(text)
    if not spans:
        return []

    max_tokens = max(1, chunk_tokens)
    safe_overlap = max(0, min(overlap_tokens, max_tokens - 1))
    step = max(1, max_tokens - safe_overlap)

    chunks: list[str] = []
    for start_index in range(0, len(spans), step):
        end_index = min(len(spans), start_index + max_tokens)
        if start_index >= end_index:
            continue
        chunk_start = spans[start_index][0]
        chunk_end = spans[end_index - 1][1]
        chunk_text = text[chunk_start:chunk_end].strip()
        if not chunk_text:
            continue
        chunks.append(chunk_text)
        if end_index >= len(spans):
            break

    return chunks


def is_heading_block(block: str) -> bool:
    normalized = block.strip()
    if not normalized:
        return False
    if MARKDOWN_HEADING_PATTERN.match(normalized):
        return True
    if "\n" in normalized:
        return False
    if NUMBERED_HEADING_PATTERN.match(normalized):
        return True

    words = normalized.split()
    if not words or len(words) > 14:
        return False
    if normalized.endswith(":"):
        return True

    letters = [char for char in normalized if char.isalpha()]
    if not letters:
        return False
    uppercase_ratio = sum(1 for char in letters if char.isupper()) / len(letters)
    return uppercase_ratio >= 0.75 and len(words) <= 10


def normalize_heading(block: str) -> str:
    normalized = block.strip()
    normalized = re.sub(r"^\s{0,3}#{1,6}\s*", "", normalized)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    return normalized.rstrip(":").strip()


def is_list_block(block: str) -> bool:
    lines = [line.strip() for line in block.splitlines() if line.strip()]
    if not lines:
        return False
    matching_lines = sum(1 for line in lines if BULLET_LINE_PATTERN.match(line))
    if matching_lines == 0:
        return False
    if matching_lines == len(lines):
        return True
    return matching_lines >= max(2, len(lines) - 1)


def clean_block(block: str) -> str:
    normalized_lines = [
        re.sub(r"[ \t]{2,}", " ", line).strip()
        for line in block.splitlines()
        if line.strip()
    ]
    return "\n".join(normalized_lines).strip()


def split_into_sections(text: str) -> list[tuple[str | None, str]]:
    raw_parts = [clean_block(part) for part in re.split(r"\n{2,}", text)]
    blocks = [part for part in raw_parts if part]
    if not blocks:
        return []

    sections: list[tuple[str | None, str]] = []
    current_title: str | None = None
    current_parts: list[str] = []

    def flush_current_section() -> None:
        if not current_parts:
            return
        body = "\n\n".join(current_parts).strip()
        if body:
            sections.append((current_title, body))
        current_parts.clear()

    for block in blocks:
        if is_heading_block(block):
            flush_current_section()
            current_title = normalize_heading(block)
            continue

        if is_list_block(block):
            flush_current_section()
            sections.append((current_title, block))
            continue

        current_parts.append(block)

    flush_current_section()

    if sections:
        return sections

    text_fallback = "\n\n".join(blocks).strip()
    if not text_fallback:
        return []
    return [(None, text_fallback)]


def render_section_text(title: str | None, body: str) -> str:
    cleaned_body = body.strip()
    if not cleaned_body:
        return title.strip() if isinstance(title, str) else ""
    if not title:
        return cleaned_body
    normalized_title = title.strip()
    if not normalized_title:
        return cleaned_body
    return f"{normalized_title}\n\n{cleaned_body}"


def merge_small_sections(
    sections: list[tuple[str | None, str]], minimum_tokens: int
) -> list[tuple[str | None, str]]:
    if not sections:
        return []

    merged: list[tuple[str | None, str]] = []
    min_tokens = max(1, minimum_tokens)

    for title, body in sections:
        normalized_body = body.strip()
        if not normalized_body:
            continue

        candidate_text = render_section_text(title, normalized_body)
        if count_tokens(candidate_text) < min_tokens and merged:
            previous_title, previous_body = merged[-1]
            combined_title = previous_title or title
            merged[-1] = (combined_title, f"{previous_body}\n\n{normalized_body}".strip())
            continue

        merged.append((title, normalized_body))

    if len(merged) >= 2:
        first_title, first_body = merged[0]
        first_text = render_section_text(first_title, first_body)
        if count_tokens(first_text) < min_tokens:
            second_title, second_body = merged[1]
            merged[1] = (
                second_title or first_title,
                f"{first_body}\n\n{second_body}".strip(),
            )
            merged = merged[1:]

    return merged


def split_section_body(
    title: str | None,
    body: str,
    max_tokens: int,
    fallback_tokens: int,
    fallback_overlap_tokens: int,
) -> list[str]:
    section_text = render_section_text(title, body)
    if not section_text:
        return []

    if count_tokens(section_text) <= max(1, max_tokens):
        return [section_text]

    if title:
        windows = split_text_with_token_windows(
            body, fallback_tokens, fallback_overlap_tokens
        )
        return [
            render_section_text(title, window)
            for window in windows
            if window and window.strip()
        ]

    return split_text_with_token_windows(
        section_text, fallback_tokens, fallback_overlap_tokens
    )


def split_documents_for_ingestion(
    documents: list[Document], content_type: str
) -> list[Document]:
    _ = content_type

    normalized_documents: list[Document] = []
    for doc in documents:
        text = normalize_chunk_text(doc.page_content)
        if not text:
            continue
        normalized_documents.append(
            Document(page_content=text, metadata=dict(doc.metadata or {}))
        )

    if not normalized_documents:
        return []

    section_chunks: list[Document] = []
    for document in normalized_documents:
        document_text = document.page_content
        base_metadata = dict(document.metadata or {})
        sections = split_into_sections(document_text)
        merged_sections = merge_small_sections(sections, settings.section_min_tokens)

        if not merged_sections:
            fallback_chunks = split_text_with_token_windows(
                document_text,
                settings.section_fallback_chunk_tokens,
                settings.section_fallback_overlap_tokens,
            )
            for chunk_text in fallback_chunks:
                section_chunks.append(
                    Document(page_content=chunk_text, metadata=dict(base_metadata))
                )
            continue

        for section_index, (section_title, section_body) in enumerate(merged_sections):
            chunks = split_section_body(
                section_title,
                section_body,
                settings.section_max_tokens,
                settings.section_fallback_chunk_tokens,
                settings.section_fallback_overlap_tokens,
            )
            for chunk_in_section, chunk_text in enumerate(chunks):
                if not chunk_text or not chunk_text.strip():
                    continue
                metadata = dict(base_metadata)
                metadata["section_index"] = section_index
                metadata["chunk_in_section"] = chunk_in_section
                section_chunks.append(
                    Document(page_content=chunk_text.strip(), metadata=metadata)
                )

    return section_chunks


def process_document(document_id: str) -> None:
    asyncio.run(_process_document(document_id))


async def _process_document(document_id: str) -> None:
    now = datetime.now(timezone.utc)
    try:
        document_object_id = ObjectId(document_id)
    except Exception:
        return

    job = get_current_job()
    job_id = job.id if job else None

    worker_client = AsyncIOMotorClient(settings.mongodb_uri)
    worker_db = worker_client[settings.database_name]
    try:
        document = await worker_db.documents.find_one({"_id": document_object_id})
        if not document:
            logger.warning(
                "Documento no encontrado", extra={"document_id": document_id}
            )
            return

        await worker_db.documents.update_one(
            {"_id": document_object_id},
            {"$set": {"status": "processing", "updated_at": now}},
        )
        if job_id:
            await worker_db.ingestion_jobs.update_one(
                {"job_id": job_id},
                {"$set": {"status": "processing", "started_at": now}},
            )

        start = time.perf_counter()
        try:
            file_path = document["file_path"]
            file_bytes = await asyncio.to_thread(download_supabase_file, file_path)
            documents = await asyncio.to_thread(
                load_documents, file_bytes, document["content_type"]
            )
            chunks = await asyncio.to_thread(
                split_documents_for_ingestion, documents, document["content_type"]
            )
            filtered_chunks = [
                chunk for chunk in chunks if chunk.page_content and chunk.page_content.strip()
            ]
            texts = [chunk.page_content for chunk in filtered_chunks]
            if not texts:
                raise RuntimeError("Documento sin texto utilizable")
            if not has_gemini_api_keys():
                raise RuntimeError("Gemini no está configurado")
            vectors = await asyncio.to_thread(
                embed_documents_with_fallback,
                texts,
                "models/gemini-embedding-001",
                settings.qdrant_vector_size,
                "RETRIEVAL_DOCUMENT",
                [document["file_name"]] * len(texts),
            )
            await upsert_chunks(document, filtered_chunks, vectors)
        except Exception as exc:
            logger.exception(
                "Ingesta fallida",
                extra={
                    "document_id": document_id,
                    "job_id": job_id,
                    "error": str(exc),
                },
            )
            await mark_failed(worker_db, document_object_id, job_id, str(exc))
            return

        latency_ms = int((time.perf_counter() - start) * 1000)
        logger.info(
            "Ingesta completada",
            extra={
                "document_id": document_id,
                "job_id": job_id,
                "chunks": len(filtered_chunks),
                "latency_ms": latency_ms,
            },
        )
        await mark_done(worker_db, document_object_id, job_id)
    finally:
        worker_client.close()


def download_supabase_file(file_path: str) -> bytes:
    if not settings.supabase_url or not settings.supabase_service_role_key:
        raise RuntimeError("Supabase no está configurado")
    if not file_path:
        raise RuntimeError("file_path inválido")

    bucket, storage_path = parse_storage_path(file_path)
    supabase_client = create_client(
        settings.supabase_url, settings.supabase_service_role_key
    )
    result = supabase_client.storage.from_(bucket).download(storage_path)
    if isinstance(result, bytes):
        return result
    return bytes(result)


def parse_storage_path(file_path: str) -> tuple[str, str]:
    if "/" in file_path:
        bucket, storage_path = file_path.split("/", 1)
        return bucket, storage_path
    if settings.supabase_storage_bucket:
        return settings.supabase_storage_bucket, file_path
    raise RuntimeError(
        "file_path debe incluir bucket/path o configurar supabase_storage_bucket"
    )


@lru_cache(maxsize=1)
def get_llamaparse_parser() -> "LlamaParse":
    if not LLAMAPARSE_AVAILABLE or LlamaParse is None:
        raise RuntimeError("LlamaParse no disponible")
    api_key = settings.llama_cloud_api_key
    if not api_key:
        raise RuntimeError("LlamaParse no esta configurado")
    return LlamaParse(
        api_key=api_key,
        result_type="markdown",
        language=settings.llamaparse_language,
        verbose=False,
    )


def load_documents_with_llamaparse(tmp_path: str, content_type: str) -> list[Document]:
    parser = get_llamaparse_parser()
    result = parser.parse(tmp_path)

    page_markdown_chunks: list[str] = []
    pages = getattr(result, "pages", None)
    if isinstance(pages, list):
        for page in pages:
            markdown = getattr(page, "md", None)
            if markdown is None and isinstance(page, dict):
                markdown = page.get("md") or page.get("text")
            if isinstance(markdown, str) and markdown.strip():
                page_markdown_chunks.append(markdown)

    if not page_markdown_chunks and hasattr(result, "get_markdown_documents"):
        try:
            markdown_documents = result.get_markdown_documents(split_by_page=True)
        except Exception:
            markdown_documents = []
        for markdown_document in markdown_documents:
            markdown = getattr(markdown_document, "text", None)
            if markdown is None and isinstance(markdown_document, dict):
                markdown = markdown_document.get("text")
            if isinstance(markdown, str) and markdown.strip():
                page_markdown_chunks.append(markdown)

    text = normalize_chunk_text("\n\n".join(page_markdown_chunks))
    if not text:
        return []
    return [
        Document(
            page_content=text,
            metadata={
                "source_parser": "llamaparse",
                "content_type": content_type,
            },
        )
    ]


def load_documents_legacy(tmp_path: str, content_type: str) -> list[Document]:
    if content_type == "pdf":
        reader = PdfReader(tmp_path)
        documents: list[Document] = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                documents.append(Document(page_content=text, metadata={"page": index}))
        return documents
    if content_type == "docx":
        text = docx2txt.process(tmp_path) or ""
        return [Document(page_content=text, metadata={})] if text.strip() else []
    if content_type == "pptx":
        presentation = Presentation(tmp_path)
        documents: list[Document] = []
        for index, slide in enumerate(presentation.slides, start=1):
            slide_sections: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    table_rows: list[str] = []
                    for row in shape.table.rows:
                        cells = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text and cell.text.strip()
                        ]
                        if cells:
                            table_rows.append(" | ".join(cells))
                    if table_rows:
                        slide_sections.append("\n".join(table_rows))
                    continue

                text = (getattr(shape, "text", "") or "").strip()
                if text:
                    slide_sections.append(text)

            if slide_sections:
                documents.append(
                    Document(
                        page_content="\n\n".join(slide_sections),
                        metadata={"page": index},
                    )
                )
        return documents

    text = Path(tmp_path).read_text(encoding="utf-8", errors="ignore")
    return [Document(page_content=text, metadata={})] if text.strip() else []


def load_documents(file_bytes: bytes, content_type: str) -> list[Document]:
    suffix_map = {"pdf": ".pdf", "docx": ".docx", "txt": ".txt", "pptx": ".pptx"}
    suffix = suffix_map.get(content_type)
    if not suffix:
        raise RuntimeError("Tipo de archivo no soportado")

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        if content_type in LLAMAPARSE_CONTENT_TYPES:
            if not settings.llama_cloud_api_key:
                logger.info(
                    "LlamaParse no esta configurado; se usa extractor legacy",
                    extra={"content_type": content_type},
                )
            elif not LLAMAPARSE_AVAILABLE:
                logger.warning(
                    "LlamaParse no disponible; se usa extractor legacy",
                    extra={"content_type": content_type},
                )
            else:
                try:
                    llamaparse_documents = load_documents_with_llamaparse(
                        tmp_path, content_type
                    )
                    if llamaparse_documents:
                        logger.info(
                            "Extraccion con LlamaParse completada",
                            extra={
                                "content_type": content_type,
                                "documents": len(llamaparse_documents),
                            },
                        )
                        return llamaparse_documents
                    logger.warning(
                        "LlamaParse no produjo texto; se usa extractor legacy",
                        extra={"content_type": content_type},
                    )
                except Exception as exc:
                    logger.warning(
                        "LlamaParse fallo; se usa extractor legacy",
                        extra={"content_type": content_type, "error": str(exc)},
                    )

        return load_documents_legacy(tmp_path, content_type)
    finally:
        try:
            Path(tmp_path).unlink(missing_ok=True)
        except FileNotFoundError:
            pass


async def upsert_chunks(
    document: dict, chunks: list[Document], vectors: list[list[float]]
) -> None:
    async with httpx.AsyncClient(base_url=settings.qdrant_url, timeout=20) as client:
        collection = settings.qdrant_collection_name
        owner_id = str(document["owner_id"])
        notebook_id = str(document["notebook_id"])
        document_id = str(document["_id"])
        file_path = document["file_path"]
        file_name = document["file_name"]
        created_at = document["created_at"].isoformat()
        source_type = document["content_type"]

        points = []
        for index, (chunk, vector) in enumerate(zip(chunks, vectors)):
            metadata = chunk.metadata or {}
            point_id = str(uuid.uuid5(uuid.NAMESPACE_URL, f"{document_id}:{index}"))
            payload = {
                "document_id": document_id,
                "notebook_id": notebook_id,
                "owner_id": owner_id,
                "file_path": file_path,
                "file_name": file_name,
                "chunk_id": index,
                "page": metadata.get("page"),
                "source_type": source_type,
                "created_at": created_at,
                "text": chunk.page_content,
            }

            points.append(
                {
                    "id": point_id,
                    "vector": vector,
                    "payload": payload,
                }
            )

        batch_size = 100
        for start in range(0, len(points), batch_size):
            batch = points[start : start + batch_size]
            response = await client.put(
                f"/collections/{collection}/points",
                json={"points": batch},
                params={"wait": "true"},
            )
            if response.status_code not in (200, 201):
                raise RuntimeError(
                    "No se pudo insertar en Qdrant: "
                    f"{response.status_code} {response.text}"
                )


async def mark_failed(
    worker_db: AsyncIOMotorDatabase,
    document_id: ObjectId,
    job_id: str | None,
    error: str,
) -> None:
    end = datetime.now(timezone.utc)
    await worker_db.documents.update_one(
        {"_id": document_id},
        {"$set": {"status": "failed", "error": error, "updated_at": end}},
    )
    if job_id:
        await worker_db.ingestion_jobs.update_one(
            {"job_id": job_id},
            {"$set": {"status": "failed", "error": error, "finished_at": end}},
        )


async def mark_done(
    worker_db: AsyncIOMotorDatabase, document_id: ObjectId, job_id: str | None
) -> None:
    end = datetime.now(timezone.utc)
    await worker_db.documents.update_one(
        {"_id": document_id},
        {"$set": {"status": "done", "error": None, "updated_at": end}},
    )
    if job_id:
        await worker_db.ingestion_jobs.update_one(
            {"job_id": job_id},
            {"$set": {"status": "done", "error": None, "finished_at": end}},
        )
